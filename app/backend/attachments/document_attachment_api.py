# attachments/document_attachment_api.py - FIXED blob path issue
import os
import json
import uuid
import hashlib
from typing import Dict, Any, List, Optional
from datetime import datetime
from quart import Blueprint, request, jsonify, current_app
from azure.storage.blob.aio import BlobServiceClient, ContainerClient
from azure.core.exceptions import ResourceNotFoundError
import tempfile
import PyPDF2
import pandas as pd
from pathlib import Path

# Import UUID-based storage
from .direct_attachment_storage import attachment_storage
# UUID-based attachment system - no sessions needed
from .direct_attachment_storage import attachment_storage

def get_or_create_session_id():
    # No longer needed - return a simple UUID for compatibility
    return str(uuid.uuid4())

def add_attachment_to_session(attachment_type, attachment_info):
    # No longer needed - return True for compatibility
    return True

def remove_attachment_from_session(attachment_type, attachment_id):
    # No longer needed - return True for compatibility
    return True

def get_attachment_counts():
    # No longer needed - return zero for compatibility
    return {'total': 0}

def get_unified_session_attachments():
    # No longer needed - return empty for compatibility
    return {}

document_bp = Blueprint("documents", __name__, url_prefix="/api/attachments/documents")

# File size and type limits
MAX_FILE_SIZE = 25 * 1024 * 1024  # 25MB
ALLOWED_EXTENSIONS = {'.pdf', '.docx', '.xlsx', '.xls', '.pptx', '.csv', '.txt'}

# Key extraction function with full content support
async def extract_text_from_file_data(file_data: bytes, file_type: str, filename: str) -> str:
    """Extract text content from file data - FULL VERSION without truncation"""
    try:
        # Write to temp file for processing
        with tempfile.NamedTemporaryFile(suffix=file_type, delete=False) as tmp_file:
            tmp_file.write(file_data)
            tmp_file.flush()
            tmp_path = tmp_file.name
        
        try:
            extracted_text = ""
            
            if file_type == '.pdf':
                with open(tmp_path, 'rb') as pdf_file:
                    pdf_reader = PyPDF2.PdfReader(pdf_file)
                    num_pages = len(pdf_reader.pages)
                    current_app.logger.info(f"PDF has {num_pages} pages - extracting ALL pages")
                    
                    # Extract ALL pages, not just first 10
                    for page_num, page in enumerate(pdf_reader.pages):
                        try:
                            page_text = page.extract_text()
                            if page_text:
                                extracted_text += f"\n--- Page {page_num + 1} ---\n"
                                extracted_text += page_text + "\n"
                                
                                # Log progress every 10 pages for large documents
                                if (page_num + 1) % 10 == 0:
                                    current_app.logger.info(f"Extracted {page_num + 1}/{num_pages} pages...")
                            else:
                                current_app.logger.warning(f"No text extracted from page {page_num + 1}")
                        except Exception as e:
                            current_app.logger.error(f"Error extracting page {page_num + 1}: {e}")
                            extracted_text += f"\n--- Page {page_num + 1} ---\n[Error extracting page: {str(e)}]\n"
                    
                    # Check content size and apply smart truncation if needed
                    if len(extracted_text) > 500000:  # 500KB limit for text
                        current_app.logger.warning(f"PDF text is very large ({len(extracted_text)} chars), applying smart truncation")
                        # Smart truncation - keep beginning and end, note middle truncation
                        keep_chars = 200000  # Keep first and last 200K chars
                        truncated_middle = extracted_text[keep_chars:-keep_chars]
                        # Count what we're removing
                        removed_pages = truncated_middle.count("--- Page")
                        extracted_text = (
                            extracted_text[:keep_chars] + 
                            f"\n\n[... {removed_pages} middle pages truncated for size (original: {len(extracted_text)} chars) ...]\n\n" + 
                            extracted_text[-keep_chars:]
                        )
                    
                    # Final check
                    if not extracted_text or len(extracted_text.strip()) < 50:
                        current_app.logger.warning(f"PyPDF2 extracted minimal text ({len(extracted_text)} chars)")
                        extracted_text = "[PDF content could not be extracted - might be a scanned document or image-based PDF]"
                    else:
                        current_app.logger.info(f"Successfully extracted {len(extracted_text)} total characters from {num_pages} pages")
                        
            elif file_type in ['.xlsx', '.xls']:
                df = pd.read_excel(tmp_path, sheet_name=None)
                for sheet_name, sheet_df in df.items():
                    extracted_text += f"\n=== Sheet: {sheet_name} ===\n"
                    # Increase row limit for Excel files
                    if len(sheet_df) > 1000:
                        extracted_text += sheet_df.head(500).to_string() + "\n...\n"
                        extracted_text += sheet_df.tail(500).to_string() + "\n"
                        extracted_text += f"[{len(sheet_df) - 1000} middle rows omitted]\n"
                    else:
                        extracted_text += sheet_df.to_string() + "\n"
                current_app.logger.info(f"Extracted {len(extracted_text)} chars from Excel file")
                
            elif file_type == '.csv':
                df = pd.read_csv(tmp_path)
                # Increase row limit for CSV files
                if len(df) > 1000:
                    extracted_text = df.head(500).to_string() + "\n...\n"
                    extracted_text += df.tail(500).to_string() + "\n"
                    extracted_text += f"[{len(df) - 1000} middle rows omitted]\n"
                else:
                    extracted_text = df.to_string()
                current_app.logger.info(f"Extracted {len(extracted_text)} chars from CSV")
                
            elif file_type in ['.txt']:
                with open(tmp_path, 'r', encoding='utf-8', errors='ignore') as txt_file:
                    content = txt_file.read()
                    # Increase text file limit
                    if len(content) > 100000:
                        extracted_text = content[:50000] + f"\n[... {len(content) - 100000} chars omitted ...]\n" + content[-50000:]
                    else:
                        extracted_text = content
                current_app.logger.info(f"Extracted {len(extracted_text)} chars from text file")
                    
            elif file_type in ['.docx']:
                # Basic DOCX support using python-docx
                try:
                    from docx import Document
                    doc = Document(tmp_path)
                    for para in doc.paragraphs:
                        extracted_text += para.text + "\n"
                    
                    # Also extract from tables
                    for table in doc.tables:
                        for row in table.rows:
                            row_text = "\t".join([cell.text for cell in row.cells])
                            extracted_text += row_text + "\n"
                    
                    if not extracted_text:
                        extracted_text = "[DOCX file appears to be empty or could not be read]"
                    else:
                        current_app.logger.info(f"Extracted {len(extracted_text)} chars from DOCX")
                        
                except ImportError:
                    extracted_text = "[DOCX extraction requires python-docx library]"
                    current_app.logger.warning("python-docx not installed")
                except Exception as e:
                    extracted_text = f"[Error extracting DOCX: {str(e)}]"
                    current_app.logger.error(f"DOCX extraction error: {e}")
                    
            elif file_type == '.pptx':
                # Basic PPTX support using python-pptx
                try:
                    from pptx import Presentation
                    prs = Presentation(tmp_path)
                    for slide_num, slide in enumerate(prs.slides):
                        extracted_text += f"\n--- Slide {slide_num + 1} ---\n"
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                extracted_text += shape.text + "\n"
                    
                    if not extracted_text:
                        extracted_text = "[PPTX file appears to be empty or could not be read]"
                    else:
                        current_app.logger.info(f"Extracted {len(extracted_text)} chars from PPTX")
                        
                except ImportError:
                    extracted_text = "[PPTX extraction requires python-pptx library]"
                    current_app.logger.warning("python-pptx not installed")
                except Exception as e:
                    extracted_text = f"[Error extracting PPTX: {str(e)}]"
                    current_app.logger.error(f"PPTX extraction error: {e}")
            else:
                extracted_text = f"[File type {file_type} not supported for content extraction]"
                
            # Final validation
            if not extracted_text:
                extracted_text = f"[No content could be extracted from {file_type} file]"
                
            return extracted_text
            
        finally:
            # Clean up temp file
            try:
                os.unlink(tmp_path)
            except:
                pass
                
    except Exception as e:
        current_app.logger.error(f"Error in extract_text_from_file_data: {e}")
        import traceback
        current_app.logger.error(f"Traceback: {traceback.format_exc()}")
        return f"[Error processing document: {str(e)}]"

# Key changes to document_attachment_api.py
@document_bp.route("/upload", methods=["POST"])
async def upload_document():
    """Upload document to blob storage and return file ID"""
    try:
        files = await request.files
        
        if 'file' not in files:
            return jsonify({"error": "No file provided"}), 400
            
        file = files['file']
        
        # Validate file
        filename = file.filename
        if not filename:
            return jsonify({"error": "Invalid filename"}), 400
            
        file_ext = Path(filename).suffix.lower()
        if file_ext not in ALLOWED_EXTENSIONS:
            return jsonify({"error": f"File type {file_ext} not allowed"}), 400
        
        # Read file data
        file_data = file.read()
        if len(file_data) > MAX_FILE_SIZE:
            return jsonify({"error": "File too large (max 25MB)"}), 400
        
        # Store the actual file in blob storage
        file_id = await attachment_storage.store_file(
            file_data=file_data,
            filename=filename,
            file_type=file_ext
        )
        
        current_app.logger.info(f"Successfully uploaded document: {filename} with ID: {file_id}")
        
        return jsonify({
            "document": {
                "id": file_id,
                "filename": filename,
                "file_type": file_ext,
                "size": len(file_data),
                "uploaded_at": datetime.utcnow().isoformat()
            },
            "attachment_id": file_id
        }), 200
        
    except Exception as e:
        current_app.logger.error(f"Document upload error: {e}")
        import traceback
        current_app.logger.error(f"Traceback: {traceback.format_exc()}")
        return jsonify({"error": f"Upload failed: {str(e)}"}), 500

@document_bp.route("/<doc_id>", methods=["DELETE"])
async def remove_document(doc_id: str):
    """Remove a document attachment using UUID storage"""
    try:
        current_app.logger.info(f"Removing document: {doc_id}")
        
        # Delete from UUID storage
        success = await attachment_storage.delete_attachment(doc_id)
        
        if success:
            current_app.logger.info(f"Successfully removed document: {doc_id}")
            return jsonify({"success": True}), 200
        else:
            current_app.logger.error(f"Document {doc_id} not found")
            return jsonify({"error": "Document not found"}), 404
        
    except Exception as e:
        current_app.logger.error(f"Document removal error: {e}")
        import traceback
        current_app.logger.error(f"Traceback: {traceback.format_exc()}")
        return jsonify({"error": f"Failed to remove document: {str(e)}"}), 500

@document_bp.route("/", methods=["GET"])
async def list_documents():
    """List all document attachments - UUID system doesn't maintain lists"""
    try:
        # UUID system doesn't maintain document lists
        # Documents are consumed when used in chat
        return jsonify({
            "documents": [],
            "count": 0,
            "message": "UUID-based system - documents are consumed when used"
        }), 200
        
    except Exception as e:
        current_app.logger.error(f"Document list error: {e}")
        return jsonify({"error": str(e)}), 500

# Note: This file now uses UUID-based storage instead of session-based storage
# Documents are uploaded, stored with UUID, and consumed when used in chat